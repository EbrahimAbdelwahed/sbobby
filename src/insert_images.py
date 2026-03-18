"""
Pipeline retroattivo: inserimento immagini da slide PPTX nelle sbobine.

Stage A1: Estrazione immagini da PPTX + dedup SHA-256
Stage A2: Pre-filtro qualità con qwen3-vl-8b via OpenRouter (con cache)
Stage A3: Selezione e placement con Gemini Flash Lite (singola chiamata per sbobina)
Stage A4: Inserimento deterministico HTML <img>

Uso standalone:
    python src/insert_images.py anatomia lezione_01 \\
        --pptx "slides/Anatomia/.../1 PANORAMICA_C2025.pptx" \\
        --pptx "slides/Anatomia/.../2 LOCOMOTORE_C2025.pptx" \\
        --output sbobine/anatomia/lezione_01_test.md
"""

import base64
import hashlib
import json
import os
import re
import shutil
import sys
import urllib.request
import urllib.error
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Optional

PROJECT_ROOT = Path(__file__).parent.parent

# ── Config ──────────────────────────────────────────────────────

OPENROUTER_URL = "https://openrouter.ai/api/v1/chat/completions"
FILTER_MODEL = "qwen/qwen3-vl-8b-instruct"
GEMINI_MODEL = "google/gemini-3.1-flash-lite-preview"

MIN_IMAGE_BYTES = 2048      # ignora icone/bullet < 2KB
MIN_IMAGE_DIM = 100         # ignora immagini < 100x100 px
MAX_IMAGES_PER_BLOCK = 3
MAX_IMAGES_SHORT_BLOCK = 1
SHORT_BLOCK_WORDS = 200
DEFAULT_WIDTH = 500


# ── Data structures ─────────────────────────────────────────────

@dataclass
class SlideImage:
    blob: bytes
    content_type: str       # e.g. "image/png"
    slide_idx: int          # 1-indexed
    slide_text: str         # concatenated text from the slide
    source_file: str        # PPTX filename
    content_hash: str = ""  # SHA-256[:16]
    filename: str = ""      # slide_001_img_1.png

    def __post_init__(self):
        if not self.content_hash:
            self.content_hash = hashlib.sha256(self.blob).hexdigest()[:16]


# ── Helpers ─────────────────────────────────────────────────────

def _load_json(path: Path) -> dict:
    if path.exists():
        return json.loads(path.read_text(encoding="utf-8"))
    return {}


def _save_json(path: Path, data: dict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")


def _cache_path(subject: str) -> Path:
    return PROJECT_ROOT / "config" / subject / "image_cache.json"


def _get_api_key() -> Optional[str]:
    key = os.environ.get("OPENROUTER_API_KEY")
    if not key:
        print("[insert_images] OPENROUTER_API_KEY non configurata")
    return key


def _ext_from_ct(ct: str) -> str:
    m = {"image/png": "png", "image/jpeg": "jpg", "image/gif": "gif",
         "image/tiff": "tiff", "image/bmp": "bmp", "image/webp": "webp"}
    return m.get(ct, "png")


def _resize_for_api(blob: bytes, content_type: str, max_side: int = 1024) -> tuple[bytes, str]:
    """Resize image if larger than max_side px."""
    try:
        from PIL import Image
        img = Image.open(BytesIO(blob))
        w, h = img.size
        if max(w, h) <= max_side:
            return blob, content_type
        ratio = max_side / max(w, h)
        new_size = (int(w * ratio), int(h * ratio))
        img = img.resize(new_size, Image.LANCZOS)
        buf = BytesIO()
        fmt = "PNG" if "png" in content_type else "JPEG"
        img.save(buf, format=fmt)
        ct = "image/png" if fmt == "PNG" else "image/jpeg"
        return buf.getvalue(), ct
    except Exception:
        return blob, content_type


def _image_dimensions(blob: bytes) -> tuple[int, int]:
    """Return (width, height) of image, or (0, 0) on error."""
    try:
        from PIL import Image
        img = Image.open(BytesIO(blob))
        return img.size
    except Exception:
        return (0, 0)


# ── A1: Estrazione immagini da PPTX ────────────────────────────

def extract_images_from_pptx(pptx_path: Path) -> list[SlideImage]:
    """Extract all images from a PPTX file with slide text and deduplication."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(pptx_path))
    images: list[SlideImage] = []
    seen_hashes: set[str] = set()

    for slide_idx, slide in enumerate(prs.slides, 1):
        # Collect slide text
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    texts.append(txt)
        slide_text = "\n".join(texts)

        # Extract images
        img_counter = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blob = shape.image.blob
                if len(blob) < MIN_IMAGE_BYTES:
                    continue

                h = hashlib.sha256(blob).hexdigest()[:16]
                if h in seen_hashes:
                    continue

                # Check minimum dimensions
                w, ht = _image_dimensions(blob)
                if w > 0 and ht > 0 and (w < MIN_IMAGE_DIM or ht < MIN_IMAGE_DIM):
                    continue

                seen_hashes.add(h)
                img_counter += 1
                ext = _ext_from_ct(shape.image.content_type)
                fname = f"slide_{slide_idx:03d}_img_{img_counter}.{ext}"

                images.append(SlideImage(
                    blob=blob,
                    content_type=shape.image.content_type,
                    slide_idx=slide_idx,
                    slide_text=slide_text,
                    source_file=pptx_path.name,
                    content_hash=h,
                    filename=fname,
                ))

    return images


def extract_all(pptx_paths: list[Path]) -> list[SlideImage]:
    """Extract images from multiple PPTX files, deduplicated globally."""
    all_images: list[SlideImage] = []
    seen: set[str] = set()

    for pptx in pptx_paths:
        print(f"[insert_images] A1: Estraggo da {pptx.name}")
        try:
            imgs = extract_images_from_pptx(pptx)
        except Exception as e:
            print(f"[insert_images] A1: ERRORE su {pptx.name}: {e}")
            continue

        for img in imgs:
            if img.content_hash not in seen:
                seen.add(img.content_hash)
                all_images.append(img)

    print(f"[insert_images] A1: {len(all_images)} immagini uniche estratte")
    return all_images


# ── A2: Pre-filtro qualità (qwen via OpenRouter) ───────────────

FILTER_PROMPT = """Classify this image. Answer ONLY with one JSON object: {"pass": true} or {"pass": false}.

KEEP (pass: true) if the image is:
- An anatomical diagram, illustration, or photograph
- A medical/clinical image (X-ray, CT, MRI, ultrasound)
- A scientific chart, graph, or schematic
- A histological or microscopic image
- A structural/functional diagram (e.g. metabolic pathway, cell structure)

DISCARD (pass: false) if the image is:
- A logo, watermark, or publisher mark
- A decorative element, background pattern, or border
- A generic clip-art or cartoon not conveying anatomical/scientific content
- A UI element, button, or navigation icon
- A photograph of a person (lecturer, stock photo) not showing anatomy
- A title slide decoration"""


def _classify_quality(image_b64: str, mime_type: str, api_key: str) -> bool:
    """Call qwen3-vl-8b on OpenRouter to classify image quality."""
    payload = json.dumps({
        "model": FILTER_MODEL,
        "messages": [
            {"role": "system", "content": FILTER_PROMPT},
            {"role": "user", "content": [
                {"type": "image_url", "image_url": {
                    "url": f"data:{mime_type};base64,{image_b64}",
                }},
            ]},
        ],
        "temperature": 0.0,
        "max_tokens": 32,
    }).encode()

    req = urllib.request.Request(
        OPENROUTER_URL, data=payload,
        headers={
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}",
        },
        method="POST",
    )

    try:
        with urllib.request.urlopen(req, timeout=30) as resp:
            data = json.loads(resp.read().decode())
        content = data["choices"][0]["message"]["content"].strip()
        # Qwen may emit <think>...</think> before JSON
        content = re.sub(r"<think>.*?</think>", "", content, flags=re.DOTALL).strip()
        result = json.loads(content)
        return bool(result.get("pass", True))
    except Exception as e:
        print(f"[insert_images] A2: Filtro errore: {e}, accetto per default")
        return True


def filter_images(images: list[SlideImage], subject: str, api_key: str) -> list[SlideImage]:
    """Filter images using qwen quality classifier with cache."""
    cache = _load_json(_cache_path(subject))

    to_classify: list[SlideImage] = []
    results: dict[str, bool] = {}

    for img in images:
        entry = cache.get(img.content_hash, {})
        if "quality_pass" in entry:
            results[img.content_hash] = entry["quality_pass"]
        else:
            to_classify.append(img)

    if to_classify:
        print(f"[insert_images] A2: Classifico {len(to_classify)} nuove immagini (10 worker)...")

        def classify_one(img: SlideImage) -> tuple[str, bool]:
            resized, ct = _resize_for_api(img.blob, img.content_type, max_side=512)
            b64 = base64.b64encode(resized).decode("ascii")
            return img.content_hash, _classify_quality(b64, ct, api_key)

        with ThreadPoolExecutor(max_workers=10) as executor:
            futures = {executor.submit(classify_one, img): img for img in to_classify}
            for future in as_completed(futures):
                h, passed = future.result()
                results[h] = passed
                entry = cache.get(h, {})
                entry["quality_pass"] = passed
                cache[h] = entry

        _save_json(_cache_path(subject), cache)
    else:
        print(f"[insert_images] A2: Tutte le immagini già classificate in cache")

    passed = [img for img in images if results.get(img.content_hash, True)]
    n_cache = len(images) - len(to_classify)
    print(f"[insert_images] A2: {len(passed)}/{len(images)} passate "
          f"({len(images) - len(passed)} scartate, {n_cache} da cache)")
    return passed


# ── A3: Selezione e placement con Gemini ────────────────────────

GEMINI_SYSTEM_PROMPT = """Sei un assistente specializzato nell'integrazione di immagini didattiche in appunti medici universitari.

Il tuo compito: data una sbobina (appunti di una lezione) e un set di immagini estratte dalle slide del docente, devi decidere QUALI immagini inserire e DOVE nel testo.

REGOLE DI SELEZIONE:
1. Inserisci SOLO immagini che aggiungono informazione visiva non derivabile dal testo (es. rapporti spaziali, morfologia 3D, imaging clinico, schemi di pathway metabolici).
2. NON inserire immagini che illustrano concetti già ben spiegati testualmente senza valore aggiunto visivo.
3. Massimo 3 immagini per blocco. Se un blocco è breve (< 200 parole), massimo 1.
4. Preferisci immagini ad alto valore didattico: radiografie reali > schemi > diagrammi generici.
5. Se più immagini mostrano lo stesso concetto, scegli la migliore e scarta le altre.

REGOLE DI POSIZIONAMENTO:
1. Le immagini vanno SEMPRE in un punto di inserimento valido (tra paragrafi, dopo elenchi, dopo tabelle).
2. MAI dentro un elenco puntato, MAI a metà di un paragrafo, MAI dentro un callout.
3. Posiziona l'immagine il più vicino possibile al testo che la descrive o la richiama.
4. Se un'immagine è rilevante per un blocco intero ma non per un punto specifico, mettila a fine blocco.

OUTPUT:
Rispondi ESCLUSIVAMENTE con un JSON array. Nessun testo prima o dopo. Nessun markdown fence.
Se nessuna immagine è rilevante per la sbobina, rispondi con [].

Ogni elemento del JSON array deve avere:
- "filename": il filename esatto dell'immagine (come indicato)
- "insertion_point_id": l'ID esatto del punto di inserimento
- "caption_it": didascalia in italiano, descrittiva e concisa (max 150 caratteri)
- "width": larghezza in pixel (default 500, usa 400 per dettagli, 600 per panoramiche)"""


def _build_gemini_message(
    sbobina_text: str,
    insertion_points: dict,
    images: list[SlideImage],
) -> list[dict]:
    """Build the multimodal user message for Gemini."""
    parts: list[dict] = []

    # Part 1: Sbobina text + insertion points
    ip_lines = []
    for ip_id, ip in sorted(insertion_points.items()):
        ip_lines.append(f"- ID: {ip_id} | Contesto: \"{ip.context}\"")

    text_content = (
        f"## SBOBINA\n\n<sbobina>\n{sbobina_text}\n</sbobina>\n\n"
        f"## PUNTI DI INSERIMENTO VALIDI\n\n<insertion_points>\n"
        + "\n".join(ip_lines)
        + "\n</insertion_points>\n\n"
        f"## IMMAGINI DISPONIBILI\n\n"
        f"Le immagini seguono in ordine. Per ogni immagine, il filename è indicato.\n"
    )
    parts.append({"type": "text", "text": text_content})

    # Part 2+: Images with labels
    for i, img in enumerate(images, 1):
        parts.append({"type": "text", "text": f"\nImmagine {i}: {img.filename}"})
        resized, ct = _resize_for_api(img.blob, img.content_type, max_side=1024)
        b64 = base64.b64encode(resized).decode("ascii")
        parts.append({
            "type": "image_url",
            "image_url": {"url": f"data:{ct};base64,{b64}"},
        })

    return parts


def _call_gemini(user_parts: list[dict], api_key: str) -> Optional[list[dict]]:
    """Call Gemini Flash Lite via OpenRouter for image placement."""
    payload = json.dumps({
        "model": GEMINI_MODEL,
        "messages": [
            {"role": "system", "content": GEMINI_SYSTEM_PROMPT},
            {"role": "user", "content": user_parts},
        ],
        "temperature": 0.1,
        "max_tokens": 4096,
    }).encode()

    req = urllib.request.Request(
        OPENROUTER_URL, data=payload,
        headers={
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}",
        },
        method="POST",
    )

    for attempt in range(3):  # max 2 retry
        try:
            with urllib.request.urlopen(req, timeout=120) as resp:
                data = json.loads(resp.read().decode())
            content = data["choices"][0]["message"]["content"].strip()
            # Strip markdown fences if present
            if content.startswith("```"):
                content = re.sub(r"^```\w*\n?", "", content)
                content = re.sub(r"\n?```$", "", content)
            result = json.loads(content)
            if isinstance(result, list):
                return result
            print(f"[insert_images] A3: Risposta non è un array, retry {attempt + 1}...")
        except json.JSONDecodeError:
            print(f"[insert_images] A3: JSON non valido, retry {attempt + 1}...")
        except urllib.error.HTTPError as e:
            body = e.read().decode() if e.fp else ""
            print(f"[insert_images] A3: HTTP {e.code}: {body[:300]}")
            return None
        except Exception as e:
            print(f"[insert_images] A3: Errore: {e}")
            if attempt < 2:
                print(f"[insert_images] A3: retry {attempt + 1}...")
            else:
                return None

    return None


def get_placements(
    sbobina_text: str,
    images: list[SlideImage],
    insertion_map: dict,
    api_key: str,
) -> list[dict]:
    """Stage A3: Call Gemini to decide which images to place where.

    Returns validated list of placement dicts:
    [{"filename": ..., "insertion_point_id": ..., "caption_it": ..., "width": ...}]
    """
    from md_parser import build_insertion_map

    if not images:
        print("[insert_images] A3: Nessuna immagine disponibile")
        return []

    # Build filename lookup
    filename_set = {img.filename for img in images}

    print(f"[insert_images] A3: Invio {len(images)} immagini + sbobina a Gemini...")
    user_parts = _build_gemini_message(sbobina_text, insertion_map, images)
    raw_placements = _call_gemini(user_parts, api_key)

    if raw_placements is None:
        print("[insert_images] A3: Gemini non ha risposto, nessun placement")
        return []

    if not raw_placements:
        print("[insert_images] A3: Gemini ha risposto con array vuoto")
        return []

    # Validate
    valid: list[dict] = []
    seen_ips: set[str] = set()

    for entry in raw_placements:
        fname = entry.get("filename", "")
        ip_id = entry.get("insertion_point_id", "")
        caption = entry.get("caption_it", "")
        width = int(entry.get("width", DEFAULT_WIDTH))

        if fname not in filename_set:
            print(f"[insert_images] A3: WARNING filename inesistente: {fname}")
            continue
        if ip_id not in insertion_map:
            print(f"[insert_images] A3: WARNING insertion_point inesistente: {ip_id}")
            continue
        if ip_id in seen_ips:
            print(f"[insert_images] A3: WARNING duplicato insertion_point: {ip_id}")
            continue

        seen_ips.add(ip_id)
        valid.append({
            "filename": fname,
            "insertion_point_id": ip_id,
            "caption_it": caption,
            "width": width,
        })

    print(f"[insert_images] A3: {len(valid)} placement validi "
          f"({len(raw_placements) - len(valid)} scartati)")
    return valid


# ── A4: Inserimento deterministico ──────────────────────────────

def _make_img_html(filename: str, caption: str, width: int, assets_rel: str) -> str:
    """Generate HTML img tag + caption for insertion."""
    # alt text: shortened caption, max 80 chars, no trailing punctuation
    alt = caption[:80].rstrip('.,;:!? ')
    return (
        f'\n<img src="{assets_rel}/{filename}" alt="{alt}" width="{width}">\n'
        f'\n*{caption}*\n'
    )


def insert_into_sbobina(
    sbobina_text: str,
    placements: list[dict],
    insertion_map: dict,
    images: list[SlideImage],
    assets_dir: Path,
    assets_rel: str,
) -> str:
    """Stage A4: Insert images at specified positions.

    Copies image files to assets_dir and inserts HTML <img> tags.
    Returns the modified sbobina text.
    """
    if not placements:
        return sbobina_text

    assets_dir.mkdir(parents=True, exist_ok=True)

    # Build image lookup by filename
    img_by_name: dict[str, SlideImage] = {img.filename: img for img in images}

    # Copy images to assets and prepare insertions
    insertions: list[tuple[int, str]] = []  # (line_number, html_to_insert)

    for p in placements:
        fname = p["filename"]
        ip_id = p["insertion_point_id"]
        caption = p["caption_it"]
        width = p["width"]

        ip = insertion_map[ip_id]
        img = img_by_name[fname]

        # Copy image to assets
        dest = assets_dir / fname
        if not dest.exists():
            dest.write_bytes(img.blob)

        html = _make_img_html(fname, caption, width, assets_rel)
        insertions.append((ip.line, html))

    # Sort insertions by line number DESCENDING (insert from bottom up)
    insertions.sort(key=lambda x: x[0], reverse=True)

    lines = sbobina_text.split('\n')
    for line_num, html in insertions:
        # Insert AFTER line_num
        html_lines = html.split('\n')
        for i, hl in enumerate(html_lines):
            lines.insert(line_num + 1 + i, hl)

    return '\n'.join(lines)


def _strip_old_markers(text: str) -> str:
    """Remove old-style image markers from previous pipeline."""
    # Remove > [immagine di: ...] placeholders
    text = re.sub(r'^> \[immagine di:\s*.+?\]\s*$', '', text, flags=re.MULTILINE)
    # Remove [CERCA_IMMAGINE: ...] markers
    text = re.sub(r'^\[CERCA_IMMAGINE:\s*.+?\]\s*$', '', text, flags=re.MULTILINE)
    # Clean up triple+ blank lines left behind
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text


# ── Orchestrator ────────────────────────────────────────────────

def _find_slides_dir(subject: str) -> Optional[Path]:
    """Auto-discover slide directory for a subject."""
    base = PROJECT_ROOT / "slides" / subject.capitalize()
    if not base.exists():
        return None
    if list(base.glob("*.pptx")):
        return base
    for sub in sorted(base.iterdir()):
        if sub.is_dir() and list(sub.glob("*.pptx")):
            return sub
    return None


def process_sbobina(
    subject: str,
    lesson: str,
    pptx_paths: Optional[list[Path]] = None,
    output_path: Optional[Path] = None,
) -> None:
    """Entry point: full retroactive pipeline A1→A4.

    Args:
        subject: e.g. "anatomia"
        lesson: e.g. "lezione_01"
        pptx_paths: explicit PPTX files to use (None = auto-discover)
        output_path: where to save result (None = overwrite sbobina)
    """
    from md_parser import build_insertion_map

    # Find sbobina
    sbobina_path = PROJECT_ROOT / "sbobine" / subject / f"{lesson}.md"
    if not sbobina_path.exists():
        raise FileNotFoundError(f"Sbobina non trovata: {sbobina_path}")

    # Find PPTX files
    if pptx_paths is None:
        slides_dir = _find_slides_dir(subject)
        if slides_dir is None:
            print(f"[insert_images] Directory slide non trovata per {subject}")
            return
        pptx_paths = sorted(slides_dir.glob("*.pptx"))

    if not pptx_paths:
        print("[insert_images] Nessun file PPTX specificato")
        return

    api_key = _get_api_key()
    if not api_key:
        return

    # Determine output path
    if output_path is None:
        output_path = sbobina_path

    # Extract lecture number for assets folder
    m = re.search(r'(\d+)', lesson)
    lecture_num = m.group(1) if m else lesson
    assets_dir = sbobina_path.parent / "assets" / f"lezione_{lecture_num}"
    assets_rel = f"assets/lezione_{lecture_num}"

    # Read sbobina and clean old markers
    sbobina_text = sbobina_path.read_text(encoding="utf-8")
    sbobina_text = _strip_old_markers(sbobina_text)

    # A1: Extract
    print(f"\n[insert_images] === A1: Estrazione immagini ===")
    images = extract_all(pptx_paths)
    if not images:
        print("[insert_images] Nessuna immagine trovata")
        return

    # A2: Quality filter
    print(f"\n[insert_images] === A2: Pre-filtro qualità ===")
    images = filter_images(images, subject, api_key)
    if not images:
        print("[insert_images] Nessuna immagine passata il filtro")
        return

    # Parse sbobina for insertion points
    print(f"\n[insert_images] === Parsing sbobina ===")
    blocks, ip_map = build_insertion_map(sbobina_text)
    print(f"[insert_images] {len(blocks)} blocchi, {len(ip_map)} punti di inserimento")

    # A3: Gemini placement
    print(f"\n[insert_images] === A3: Selezione e placement (Gemini) ===")
    placements = get_placements(sbobina_text, images, ip_map, api_key)
    if not placements:
        print("[insert_images] Nessun placement, sbobina invariata")
        # Still save cleaned version (old markers removed)
        output_path.write_text(sbobina_text, encoding="utf-8")
        return

    # A4: Insert
    print(f"\n[insert_images] === A4: Inserimento deterministico ===")
    result = insert_into_sbobina(
        sbobina_text, placements, ip_map, images, assets_dir, assets_rel,
    )

    # Save
    output_path.write_text(result, encoding="utf-8")

    # Report
    n_images = len(set(p["filename"] for p in placements))
    n_blocks = len(set(
        p["insertion_point_id"].split(".")[0] for p in placements
    ))
    print(f"\n[insert_images] === Report ===")
    print(f"  Immagini inserite: {n_images}")
    print(f"  Blocchi con immagini: {n_blocks}/{len(blocks)}")
    print(f"  Immagini candidate: {len(images)} (post-filtro)")
    print(f"  Assets: {assets_dir}")
    print(f"  Output: {output_path}")


# ── CLI ─────────────────────────────────────────────────────────

def main():
    import argparse
    parser = argparse.ArgumentParser(description="Inserimento immagini nelle sbobine")
    parser.add_argument("subject", help="Materia (es. anatomia)")
    parser.add_argument("lesson", help="Lezione (es. lezione_01)")
    parser.add_argument("--pptx", action="append", dest="pptx_files",
                        help="File PPTX da usare (ripetibile)")
    parser.add_argument("--output", "-o", help="File output (default: sovrascrive)")

    args = parser.parse_args()

    pptx_paths = [Path(p) for p in args.pptx_files] if args.pptx_files else None
    output_path = Path(args.output) if args.output else None

    process_sbobina(args.subject, args.lesson, pptx_paths, output_path)


if __name__ == "__main__":
    main()
