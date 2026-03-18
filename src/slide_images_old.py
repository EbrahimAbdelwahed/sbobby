"""
Estrazione immagini dalle diapositive, captioning VL e inserimento in sbobina.

Uso standalone:
    python src/slide_images.py anatomia lezione_01
    python src/slide_images.py anatomia lezione_01 --slides-dir slides/Anatomia/cartella/

Fasi:
    1. Estrae immagini singole + testo da PPTX/PDF
    2. Calcola hash per cache (riuso tra lezioni con stesse slide)
    3. Caption via Gemini Flash Lite su OpenRouter (immagine + testo slide)
    4. LLM legge sbobina e suggerisce punti di inserimento [CERCA_IMMAGINE: ...]
    5. Embedding match (OpenRouter) tra marker e caption → inserimento con posizionamento automatico
"""

import base64
import hashlib
import json
import os
import re
import sys
import urllib.request
import urllib.error
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Optional

PROJECT_ROOT = Path(__file__).parent.parent

# Dimensione minima immagine in byte (ignora icone/bullet decorativi)
MIN_IMAGE_BYTES = 2048

# Modelli VL su OpenRouter
VL_MODEL = "google/gemini-3.1-flash-lite-preview"
FILTER_MODEL = "qwen/qwen3-vl-8b-instruct"
OPENROUTER_URL = "https://openrouter.ai/api/v1/chat/completions"

# Percorsi prompt
FILTER_PROMPT_PATH = PROJECT_ROOT / "prompts" / "image_filter.md"
CAPTION_PROMPT_PATH = PROJECT_ROOT / "prompts" / "image_caption.md"
REVIEWER_PROMPT_PATH = PROJECT_ROOT / "prompts" / "insert_images.md"


# ── Data structures ──────────────────────────────────────────────

@dataclass
class SlideImage:
    blob: bytes
    content_type: str  # e.g. "image/png"
    slide_number: int
    slide_text: str
    source_file: str
    content_hash: str = field(default="")

    def __post_init__(self):
        if not self.content_hash:
            self.content_hash = hashlib.sha256(self.blob).hexdigest()[:16]


@dataclass
class CaptionedImage:
    content_hash: str
    filename: str          # senza estensione
    brief_caption: str     # alt-text breve per Obsidian
    caption_embedding: str # descrizione ricca per matching semantico
    content_type: str
    slide_number: int
    slide_text: str
    source_file: str
    placement: str = "side"   # "side" | "full-width"
    width_px: int = 250        # rilevante solo per placement="side"
    saved_path: Optional[Path] = None


# ── Cache ────────────────────────────────────────────────────────

def _load_json(path: Path) -> dict:
    if path.exists():
        return json.loads(path.read_text(encoding="utf-8"))
    return {}


def _save_json(path: Path, data: dict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")


def _cache_path(subject: str) -> Path:
    return PROJECT_ROOT / "config" / subject / "image_cache.json"


# ── Extraction: PPTX ────────────────────────────────────────────

def extract_images_from_pptx(pptx_path: Path) -> list[SlideImage]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(pptx_path))
    images = []
    seen_hashes = set()

    for slide_idx, slide in enumerate(prs.slides, 1):
        # Raccogli testo della slide
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    texts.append(txt)
        slide_text = "\n".join(texts)

        # Estrai immagini
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blob = shape.image.blob
                if len(blob) < MIN_IMAGE_BYTES:
                    continue
                h = hashlib.sha256(blob).hexdigest()[:16]
                if h in seen_hashes:
                    continue
                seen_hashes.add(h)

                images.append(SlideImage(
                    blob=blob,
                    content_type=shape.image.content_type,
                    slide_number=slide_idx,
                    slide_text=slide_text,
                    source_file=pptx_path.name,
                    content_hash=h,
                ))

    return images


# ── Extraction: PDF ──────────────────────────────────────────────

def extract_images_from_pdf(pdf_path: Path) -> list[SlideImage]:
    import fitz

    doc = fitz.open(str(pdf_path))
    images = []
    seen_hashes = set()

    for page_idx in range(len(doc)):
        page = doc[page_idx]
        page_text = page.get_text().strip()

        for img_info in page.get_images(full=True):
            xref = img_info[0]
            try:
                img_data = doc.extract_image(xref)
            except Exception:
                continue
            blob = img_data["image"]
            if len(blob) < MIN_IMAGE_BYTES:
                continue

            h = hashlib.sha256(blob).hexdigest()[:16]
            if h in seen_hashes:
                continue
            seen_hashes.add(h)

            ext = img_data.get("ext", "png")
            ct = f"image/{ext}" if ext != "jpg" else "image/jpeg"

            images.append(SlideImage(
                blob=blob,
                content_type=ct,
                slide_number=page_idx + 1,
                slide_text=page_text,
                source_file=pdf_path.name,
                content_hash=h,
            ))

    doc.close()
    return images


# ── Extract all from directory ───────────────────────────────────

def _find_slides_dir(subject: str) -> Optional[Path]:
    """Trova la cartella slide per la materia. Cerca in slides/{Subject}/."""
    base = PROJECT_ROOT / "slides" / subject.capitalize()
    if not base.exists():
        return None
    # Se la cartella contiene direttamente file pptx/pdf, usala
    if list(base.glob("*.pptx")) or list(base.glob("*.pdf")):
        return base
    # Altrimenti cerca la prima sottocartella che contiene slide
    for sub in sorted(base.iterdir()):
        if sub.is_dir() and (list(sub.glob("*.pptx")) or list(sub.glob("*.pdf"))):
            return sub
    return None


def extract_all_slides(slides_dir: Path) -> list[SlideImage]:
    """Estrae tutte le immagini da PPTX e PDF nella directory, deduplicate per hash."""
    all_images = []
    seen = set()

    files = sorted(list(slides_dir.glob("*.pptx")) + list(slides_dir.glob("*.pdf")))
    if not files:
        print(f"[slide_images] Nessun file PPTX/PDF trovato in {slides_dir}")
        return []

    print(f"[slide_images] Trovati {len(files)} file di slide")

    for f in files:
        print(f"[slide_images] Estraggo da: {f.name}")
        try:
            if f.suffix.lower() == ".pptx":
                imgs = extract_images_from_pptx(f)
            else:
                imgs = extract_images_from_pdf(f)
        except Exception as e:
            print(f"[slide_images] ERRORE su {f.name}: {e}, skip")
            continue

        for img in imgs:
            if img.content_hash not in seen:
                seen.add(img.content_hash)
                all_images.append(img)

    print(f"[slide_images] Totale immagini uniche estratte: {len(all_images)}")
    return all_images


# ── OpenRouter VL API ────────────────────────────────────────────

def _get_openrouter_key() -> Optional[str]:
    key = os.environ.get("OPENROUTER_API_KEY")
    if not key:
        print("[slide_images] OPENROUTER_API_KEY non configurata, skip captioning")
    return key


def _resize_for_api(blob: bytes, content_type: str, max_side: int = 1024) -> tuple[bytes, str]:
    """Ridimensiona l'immagine se più grande di max_side px. Ritorna (blob, content_type)."""
    try:
        from PIL import Image
        img = Image.open(BytesIO(blob))
        w, h = img.size
        if max(w, h) <= max_side:
            return blob, content_type
        ratio = max_side / max(w, h)
        new_size = (int(w * ratio), int(h * ratio))
        img = img.resize(new_size, Image.LANCZOS)
        buf = BytesIO()
        fmt = "PNG" if "png" in content_type else "JPEG"
        img.save(buf, format=fmt)
        ct = "image/png" if fmt == "PNG" else "image/jpeg"
        return buf.getvalue(), ct
    except Exception:
        return blob, content_type


def _call_openrouter_vl(
    image_b64: str,
    mime_type: str,
    slide_text: str,
    subject: str,
    api_key: str,
) -> Optional[dict]:
    """Chiama OpenRouter VL API e restituisce dict con filename/brief_caption/caption_embedding/placement/width_px."""
    caption_prompt = CAPTION_PROMPT_PATH.read_text(encoding="utf-8")
    user_text = f"Materia: {subject}\nTesto sulla slide: {slide_text}" if slide_text else f"Materia: {subject}"

    payload = json.dumps({
        "model": VL_MODEL,
        "messages": [
            {"role": "system", "content": caption_prompt},
            {"role": "user", "content": [
                {"type": "image_url", "image_url": {
                    "url": f"data:{mime_type};base64,{image_b64}",
                }},
                {"type": "text", "text": user_text},
            ]},
        ],
        "temperature": 0.2,
        "max_tokens": 512,
    }).encode()

    req = urllib.request.Request(
        OPENROUTER_URL,
        data=payload,
        headers={
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}",
        },
        method="POST",
    )

    for attempt in range(2):
        try:
            with urllib.request.urlopen(req, timeout=60) as resp:
                data = json.loads(resp.read().decode())
            content = data["choices"][0]["message"]["content"]
            # Rimuovi eventuale fencing markdown
            content = content.strip()
            if content.startswith("```"):
                content = re.sub(r"^```\w*\n?", "", content)
                content = re.sub(r"\n?```$", "", content)
            result = json.loads(content)
            if "filename" in result and "brief_caption" in result:
                return result
            if attempt == 0:
                print("[slide_images] JSON incompleto dal VL, retry...")
        except json.JSONDecodeError:
            if attempt == 0:
                print("[slide_images] JSON non valido dal VL, retry...")
        except urllib.error.HTTPError as e:
            body = e.read().decode() if e.fp else ""
            print(f"[slide_images] OpenRouter HTTP {e.code}: {body[:200]}")
            return None
        except Exception as e:
            print(f"[slide_images] Errore OpenRouter: {e}")
            return None

    return None


# ── Pre-filtro qualità immagini ──────────────────────────────────

def _classify_image_quality(image_b64: str, mime_type: str, api_key: str) -> bool:
    """Usa qwen3-vl-8b per determinare se l'immagine è accademica/medica adeguata."""
    filter_prompt = FILTER_PROMPT_PATH.read_text(encoding="utf-8")
    payload = json.dumps({
        "model": FILTER_MODEL,
        "messages": [
            {"role": "system", "content": filter_prompt},
            {"role": "user", "content": [
                {"type": "image_url", "image_url": {
                    "url": f"data:{mime_type};base64,{image_b64}",
                }},
            ]},
        ],
        "temperature": 0.0,
        "max_tokens": 32,
    }).encode()

    req = urllib.request.Request(
        OPENROUTER_URL,
        data=payload,
        headers={
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}",
        },
        method="POST",
    )

    try:
        with urllib.request.urlopen(req, timeout=30) as resp:
            data = json.loads(resp.read().decode())
        content = data["choices"][0]["message"]["content"].strip()
        # Qwen può emettere tag <think>...</think> prima del JSON
        content = re.sub(r"<think>.*?</think>", "", content, flags=re.DOTALL).strip()
        result = json.loads(content)
        return bool(result.get("pass", True))
    except Exception as e:
        print(f"[slide_images] Filtro qualità errore: {e}, accetto per default")
        return True  # fallback conservativo: non scartare se errore


def _filter_images_parallel(
    slide_images: list[SlideImage],
    subject: str,
    api_key: str,
) -> list[SlideImage]:
    """Filtra immagini non accademiche con qwen3-vl-8b (10 worker paralleli).
    Salva il risultato del filtro in image_cache.json per evitare ri-classificazioni.
    """
    cache = _load_json(_cache_path(subject))

    to_classify: list[SlideImage] = []
    results_map: dict[str, bool] = {}

    for img in slide_images:
        entry = cache.get(img.content_hash, {})
        if "quality_pass" in entry:
            results_map[img.content_hash] = entry["quality_pass"]
        else:
            to_classify.append(img)

    if to_classify:
        print(f"[slide_images] Filtro qualità: classifico {len(to_classify)} nuove immagini "
              f"(10 worker)...")

        def classify_one(img: SlideImage) -> tuple[str, bool]:
            resized, ct = _resize_for_api(img.blob, img.content_type, max_side=512)
            b64 = base64.b64encode(resized).decode("ascii")
            passed = _classify_image_quality(b64, ct, api_key)
            return img.content_hash, passed

        with ThreadPoolExecutor(max_workers=10) as executor:
            futures = {executor.submit(classify_one, img): img for img in to_classify}
            done = 0
            for future in as_completed(futures):
                done += 1
                h, passed = future.result()
                results_map[h] = passed
                entry = cache.get(h, {})
                entry["quality_pass"] = passed
                cache[h] = entry

        _save_json(_cache_path(subject), cache)

    passed_imgs = [img for img in slide_images if results_map.get(img.content_hash, True)]
    n_failed = len(slide_images) - len(passed_imgs)
    n_from_cache = len(slide_images) - len(to_classify)
    print(f"[slide_images] Filtro qualità: {len(passed_imgs)}/{len(slide_images)} passate "
          f"({n_failed} scartate, {n_from_cache} da cache)")
    return passed_imgs


# ── Captioning ───────────────────────────────────────────────────

def _ext_from_content_type(ct: str) -> str:
    mapping = {"image/png": "png", "image/jpeg": "jpg", "image/gif": "gif",
               "image/tiff": "tiff", "image/bmp": "bmp", "image/webp": "webp"}
    return mapping.get(ct, "png")


def _sanitize_filename(name: str) -> str:
    """Normalizza il filename: snake_case, solo lettere/numeri/underscore, max 50 chars."""
    name = name.lower().strip()
    name = re.sub(r"[^a-z0-9_àèéìòùç]", "_", name)
    name = re.sub(r"_+", "_", name).strip("_")
    return name[:50]


def _unique_filename(name: str, existing: set[str]) -> str:
    """Aggiunge suffisso numerico se il nome esiste già."""
    if name not in existing:
        return name
    i = 2
    while f"{name}_{i}" in existing:
        i += 1
    return f"{name}_{i}"


def caption_images(
    slide_images: list[SlideImage],
    subject: str,
    output_dir: Path,
) -> list[CaptionedImage]:
    """Genera caption per tutte le immagini (10 worker paralleli per le nuove), usando cache."""
    cache = _load_json(_cache_path(subject))
    api_key = _get_openrouter_key()
    output_dir.mkdir(parents=True, exist_ok=True)

    # Separa immagini già captionate (hanno "filename") da quelle nuove
    # Nota: _filter_images_parallel scrive solo "quality_pass" — non basta per essere "cached"
    cached_hashes = {
        img.content_hash for img in slide_images
        if "filename" in cache.get(img.content_hash, {})
    }
    new_imgs = [img for img in slide_images if img.content_hash not in cached_hashes]

    # Chiama API in parallelo per le nuove immagini
    api_results: dict[str, Optional[dict]] = {}
    if new_imgs and api_key:
        print(f"[slide_images] Captioning {len(new_imgs)} nuove immagini (10 worker)...")

        def caption_one(img: SlideImage) -> tuple[str, Optional[dict]]:
            resized_blob, resized_ct = _resize_for_api(img.blob, img.content_type)
            b64 = base64.b64encode(resized_blob).decode("ascii")
            result = _call_openrouter_vl(b64, resized_ct, img.slide_text, subject, api_key)
            return img.content_hash, result

        with ThreadPoolExecutor(max_workers=10) as executor:
            futures = {executor.submit(caption_one, img): img for img in new_imgs}
            done = 0
            for future in as_completed(futures):
                done += 1
                h, result = future.result()
                api_results[h] = result
                print(f"[slide_images]   {done}/{len(new_imgs)} completate")

    # Costruisci risultati in ordine (sequenziale per gestire filename univoci)
    results: list[CaptionedImage] = []
    used_filenames: set[str] = set()
    cached_count = 0
    captioned_count = 0

    for img in slide_images:
        ext = _ext_from_content_type(img.content_type)

        if img.content_hash in cached_hashes:
            entry = cache[img.content_hash]
            fname = _unique_filename(entry["filename"], used_filenames)
            used_filenames.add(fname)
            saved_path = output_dir / f"{fname}.{ext}"
            if not saved_path.exists():
                saved_path.write_bytes(img.blob)
            # Compatibilità backward: cache vecchia usa rich_caption
            cap_emb = entry.get("caption_embedding", entry.get("rich_caption", ""))
            results.append(CaptionedImage(
                content_hash=img.content_hash,
                filename=fname,
                brief_caption=entry.get("brief_caption", ""),
                caption_embedding=cap_emb,
                content_type=img.content_type,
                slide_number=img.slide_number,
                slide_text=img.slide_text,
                source_file=img.source_file,
                placement=entry.get("placement", "side"),
                width_px=int(entry.get("width_px", 250)),
                saved_path=saved_path,
            ))
            cached_count += 1
            continue

        # Nuova immagine
        vl_result = api_results.get(img.content_hash) if api_key else None

        if vl_result:
            fname = _sanitize_filename(vl_result["filename"])
            brief = vl_result.get("brief_caption", "")
            cap_emb = vl_result.get("caption_embedding", "")
            placement = vl_result.get("placement", "side")
            width_px = int(vl_result.get("width_px", 250))
        else:
            fname = img.content_hash
            brief = img.slide_text[:120] if img.slide_text else ""
            cap_emb = ""
            placement = "side"
            width_px = 250

        fname = _unique_filename(fname, used_filenames)
        used_filenames.add(fname)

        saved_path = output_dir / f"{fname}.{ext}"
        saved_path.write_bytes(img.blob)

        cache[img.content_hash] = {
            "filename": fname,
            "brief_caption": brief,
            "caption_embedding": cap_emb,
            "placement": placement,
            "width_px": width_px,
            "quality_pass": True,
            "source_file": img.source_file,
            "slide_number": img.slide_number,
        }
        captioned_count += 1

        results.append(CaptionedImage(
            content_hash=img.content_hash,
            filename=fname,
            brief_caption=brief,
            caption_embedding=cap_emb,
            content_type=img.content_type,
            slide_number=img.slide_number,
            slide_text=img.slide_text,
            source_file=img.source_file,
            placement=placement,
            width_px=width_px,
            saved_path=saved_path,
        ))

    _save_json(_cache_path(subject), cache)
    print(f"[slide_images] Caption: {captioned_count} nuove, {cached_count} da cache")
    return results


# ── Embedding matching ───────────────────────────────────────────

MATCH_THRESHOLD = 0.65  # soglia minima cosine similarity per accettare un match
EMBEDDING_MODEL = "openai/text-embedding-3-large"
EMBEDDING_URL = "https://openrouter.ai/api/v1/embeddings"
EMBEDDING_BATCH_SIZE = 100  # max testi per singola chiamata API


def _embedding_cache_path(subject: str) -> Path:
    return PROJECT_ROOT / "config" / subject / "image_embeddings.json"


def _call_embeddings_api(texts: list[str], api_key: str) -> list[list[float]]:
    """Chiama OpenRouter embeddings API. Ritorna lista di vettori."""
    results: list[list[float]] = []
    for i in range(0, len(texts), EMBEDDING_BATCH_SIZE):
        batch = texts[i:i + EMBEDDING_BATCH_SIZE]
        payload = json.dumps({
            "model": EMBEDDING_MODEL,
            "input": batch,
        }).encode()

        req = urllib.request.Request(
            EMBEDDING_URL,
            data=payload,
            headers={
                "Content-Type": "application/json",
                "Authorization": f"Bearer {api_key}",
            },
            method="POST",
        )

        try:
            with urllib.request.urlopen(req, timeout=60) as resp:
                data = json.loads(resp.read().decode())
            # Ordina per index (l'API potrebbe non rispettare l'ordine)
            sorted_data = sorted(data["data"], key=lambda x: x["index"])
            results.extend([item["embedding"] for item in sorted_data])
        except Exception as e:
            print(f"[slide_images] Errore embeddings API: {e}")
            # Ritorna vettori zero per questo batch
            results.extend([[0.0] for _ in batch])

    return results


def _cosine_similarity(a: list[float], b: list[float]) -> float:
    """Cosine similarity tra due vettori."""
    dot = sum(x * y for x, y in zip(a, b))
    norm_a = sum(x * x for x in a) ** 0.5
    norm_b = sum(x * x for x in b) ** 0.5
    if norm_a == 0 or norm_b == 0:
        return 0.0
    return dot / (norm_a * norm_b)


def _build_caption_embeddings(
    captioned_images: list[CaptionedImage],
    subject: str,
    api_key: str,
) -> dict[str, list[float]]:
    """Calcola embeddings per tutte le caption, con cache su disco."""
    cache_path = _embedding_cache_path(subject)
    cache = _load_json(cache_path)

    # Trova immagini senza embedding in cache
    to_embed: list[tuple[str, str]] = []  # (content_hash, testo)
    for img in captioned_images:
        if img.content_hash not in cache:
            # caption_embedding è il campo ottimizzato per matching; brief come fallback
            text = (img.caption_embedding or img.brief_caption).strip()
            if text:
                to_embed.append((img.content_hash, text))

    if to_embed:
        print(f"[slide_images] Embedding {len(to_embed)} caption nuove...")
        texts = [t for _, t in to_embed]
        vectors = _call_embeddings_api(texts, api_key)
        for (h, _), vec in zip(to_embed, vectors):
            cache[h] = vec
        _save_json(cache_path, cache)
        print(f"[slide_images] Embedding salvati in cache ({len(cache)} totali)")
    else:
        print(f"[slide_images] Embedding caption: tutte da cache ({len(cache)})")

    return cache


def _embed_match(
    query_embedding: list[float],
    captioned_images: list[CaptionedImage],
    caption_embeddings: dict[str, list[float]],
    used: set[str],
) -> tuple[Optional[CaptionedImage], float]:
    """Trova la migliore immagine per una query via cosine similarity."""
    best_score = 0.0
    best_img = None

    for img in captioned_images:
        if img.content_hash in used:
            continue
        emb = caption_embeddings.get(img.content_hash)
        if not emb or emb == [0.0]:
            continue
        score = _cosine_similarity(query_embedding, emb)
        if score > best_score:
            best_score = score
            best_img = img

    return best_img, best_score


# ── Posizionamento e dimensione ──────────────────────────────────


def _format_embed(img: CaptionedImage) -> str:
    """Genera il markup Obsidian usando placement e width_px decisi da Gemini."""
    ext = _ext_from_content_type(img.content_type)
    full_name = f"{img.filename}.{ext}"

    if img.placement == "full-width":
        return f"![[{full_name}]]"
    else:
        width = img.width_px if img.width_px else 250
        return f"> [!img-right]\n> ![[{full_name}|{width}]]"


# ── Splitting sbobina in blocchi ──────────────────────────────────

CERCA_RE = re.compile(r"^\[CERCA_IMMAGINE:\s*(.+?)\]\s*$", re.MULTILINE)
PLACEHOLDER_RE = re.compile(r"^> \[immagine di:\s*(.+?)\]\s*$", re.MULTILINE)
BLOCK_RE = re.compile(r"^## BLOCCO\b", re.MULTILINE)


def _split_into_blocks(sbobina_text: str) -> tuple[str, list[str], str]:
    """Splitta la sbobina in header, blocchi, e footer (appendice + chiusura).

    Ritorna (header, [blocco1, blocco2, ...], footer).
    Header = tutto prima del primo ## BLOCCO.
    Footer = dall'APPENDICE TABELLARE in poi (o stringa vuota).
    """
    # Trova tutte le posizioni dei ## BLOCCO
    block_starts = [m.start() for m in BLOCK_RE.finditer(sbobina_text)]

    if not block_starts:
        # Nessun blocco trovato — ritorna tutto come header
        return sbobina_text, [], ""

    header = sbobina_text[:block_starts[0]]

    # Trova appendice/chiusura
    appendice_match = re.search(r"^## APPENDICE\b", sbobina_text, re.MULTILINE)
    footer_start = appendice_match.start() if appendice_match else len(sbobina_text)

    # Estrai blocchi individuali
    blocks = []
    for i, start in enumerate(block_starts):
        if start >= footer_start:
            break
        end = block_starts[i + 1] if i + 1 < len(block_starts) else footer_start
        end = min(end, footer_start)
        blocks.append(sbobina_text[start:end])

    footer = sbobina_text[footer_start:]

    return header, blocks, footer


# ── Step 1: LLM suggerisce punti di inserimento (blocco per blocco) ──

def _llm_suggest_for_block(block_text: str, block_idx: int, total: int) -> str:
    """LLM legge un singolo blocco e aggiunge marker [CERCA_IMMAGINE: ...]."""
    from api_client import chat

    system_prompt = REVIEWER_PROMPT_PATH.read_text(encoding="utf-8")

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": block_text},
    ]

    print(f"[slide_images] LLM blocco {block_idx}/{total}...")
    result = chat(messages, temperature=0.3)  # DeepSeek chat, 8K output basta per un blocco

    # Validazione: l'output deve contenere il titolo del blocco
    first_line = block_text.strip().split("\n")[0]
    if first_line not in result:
        print(f"[slide_images]   ATTENZIONE: blocco {block_idx} alterato, uso originale")
        return block_text

    return result


def _llm_suggest_insertion_points(sbobina_text: str) -> str:
    """Processa la sbobina blocco per blocco, aggiungendo marker [CERCA_IMMAGINE: ...]."""
    header, blocks, footer = _split_into_blocks(sbobina_text)

    if not blocks:
        print("[slide_images] Nessun blocco trovato nella sbobina, skip suggerimenti LLM")
        return sbobina_text

    print(f"[slide_images] Sbobina splittata: {len(blocks)} blocchi")

    annotated_blocks = []
    for i, block in enumerate(blocks, 1):
        annotated = _llm_suggest_for_block(block, i, len(blocks))
        annotated_blocks.append(annotated)

    return header + "".join(annotated_blocks) + footer


# ── Step 2: Embedding match e sostituzione ───────────────────────

def _resolve_markers(
    text: str,
    captioned_images: list[CaptionedImage],
    subject: str,
) -> str:
    """Sostituisce [CERCA_IMMAGINE: ...] e > [immagine di: ...] con embed reali via embeddings."""
    api_key = _get_openrouter_key()
    if not api_key:
        print("[slide_images] OPENROUTER_API_KEY mancante, skip matching")
        return text

    # 1. Costruisci/carica embeddings delle caption
    caption_embeddings = _build_caption_embeddings(captioned_images, subject, api_key)

    # 2. Raccogli tutte le query dai marker
    placeholder_matches = list(PLACEHOLDER_RE.finditer(text))
    cerca_matches = list(CERCA_RE.finditer(text))
    all_queries = [m.group(1) for m in placeholder_matches] + [m.group(1) for m in cerca_matches]

    if not all_queries:
        print("[slide_images] Nessun marker da risolvere")
        return text

    # 3. Batch embed di tutte le query
    print(f"[slide_images] Embedding {len(all_queries)} query...")
    query_embeddings = _call_embeddings_api(all_queries, api_key)

    # Mappa query → embedding
    query_emb_map: dict[str, list[float]] = {}
    for query, emb in zip(all_queries, query_embeddings):
        query_emb_map[query] = emb

    # 4. Risolvi marker
    used: set[str] = set()
    match_log: list[str] = []

    def replace_cerca(m):
        query = m.group(1)
        emb = query_emb_map.get(query, [0.0])
        img, score = _embed_match(emb, captioned_images, caption_embeddings, used)
        if img and score >= MATCH_THRESHOLD:
            used.add(img.content_hash)
            embed = _format_embed(img)
            match_log.append(f"  MATCH ({score:.2f}): '{query[:60]}' → {img.filename}")
            return embed
        match_log.append(f"  SKIP  ({score:.2f}): '{query[:60]}' — sotto soglia")
        return ""  # rimuove marker senza match

    def replace_placeholder(m):
        query = m.group(1)
        emb = query_emb_map.get(query, [0.0])
        img, score = _embed_match(emb, captioned_images, caption_embeddings, used)
        if img and score >= MATCH_THRESHOLD:
            used.add(img.content_hash)
            embed = _format_embed(img)
            match_log.append(f"  MATCH ({score:.2f}): placeholder '{query[:60]}' → {img.filename}")
            return embed
        match_log.append(f"  KEEP  ({score:.2f}): placeholder '{query[:60]}' — nessun match")
        return m.group(0)  # lascia il placeholder invariato

    # Prima i placeholder esistenti (priorità), poi i marker LLM
    text = PLACEHOLDER_RE.sub(replace_placeholder, text)
    text = CERCA_RE.sub(replace_cerca, text)

    # Pulisci righe vuote multiple lasciate da marker rimossi
    text = re.sub(r"\n{3,}", "\n\n", text)

    # Log match
    log_path = PROJECT_ROOT / "config" / subject / "image_matches.log"
    log_path.write_text(
        f"Soglia: {MATCH_THRESHOLD}\n"
        f"Immagini inserite: {len(used)}/{len(captioned_images)}\n\n"
        + "\n".join(match_log),
        encoding="utf-8",
    )
    print(f"[slide_images] Match: {len(used)} immagini inserite "
          f"(soglia {MATCH_THRESHOLD})")
    if match_log:
        for line in match_log:
            print(f"[slide_images] {line}")

    return text


# ── Entry point ──────────────────────────────────────────────────

def process_sbobina(
    subject: str,
    lesson: str,
    slides_dir: Optional[Path] = None,
) -> None:
    """Entry point: estrae, captiona e inserisce immagini nella sbobina."""
    # Trova sbobina
    sbobina_path = PROJECT_ROOT / "sbobine" / subject / f"{lesson}.md"
    if not sbobina_path.exists():
        raise FileNotFoundError(f"Sbobina non trovata: {sbobina_path}")

    # Trova directory slide
    if slides_dir is None:
        slides_dir = _find_slides_dir(subject)
    if slides_dir is None or not slides_dir.exists():
        print(f"[slide_images] Directory slide non trovata per {subject}. "
              f"Atteso: slides/{subject.capitalize()}/")
        return

    # Directory output immagini
    img_output_dir = PROJECT_ROOT / "sbobine" / "strutture"

    # 1. Estrai immagini
    print(f"\n[slide_images] === Estrazione immagini da {slides_dir.name} ===")
    slide_imgs = extract_all_slides(slides_dir)
    if not slide_imgs:
        print("[slide_images] Nessuna immagine trovata, skip")
        return

    # 2. Pre-filtro qualità (qwen3-vl-8b, parallelo)
    api_key = _get_openrouter_key()
    if api_key:
        print(f"\n[slide_images] === Pre-filtro qualità immagini ===")
        slide_imgs = _filter_images_parallel(slide_imgs, subject, api_key)
        if not slide_imgs:
            print("[slide_images] Nessuna immagine passata il filtro, skip")
            return

    # 3. Caption (Gemini, 10 worker paralleli, usa cache se disponibile)
    print(f"\n[slide_images] === Captioning {len(slide_imgs)} immagini ===")
    captioned = caption_images(slide_imgs, subject, img_output_dir)
    if not captioned:
        print("[slide_images] Nessuna immagine processata, skip")
        return

    # 4. LLM suggerisce punti di inserimento (blocco per blocco)
    sbobina_text = sbobina_path.read_text(encoding="utf-8")
    print(f"\n[slide_images] === Suggerimento punti di inserimento ===")
    annotated = _llm_suggest_insertion_points(sbobina_text)

    # Conta marker trovati
    cerca_count = len(CERCA_RE.findall(annotated))
    placeholder_count = len(PLACEHOLDER_RE.findall(annotated))
    print(f"[slide_images] Marker trovati: {cerca_count} [CERCA_IMMAGINE] + "
          f"{placeholder_count} [immagine di]")

    # 5. Embedding match e inserimento
    print(f"\n[slide_images] === Matching e inserimento ===")
    result = _resolve_markers(annotated, captioned, subject)

    # Salva
    sbobina_path.write_text(result, encoding="utf-8")
    print(f"[slide_images] Sbobina aggiornata: {sbobina_path}")
    print(f"[slide_images] Immagini in: {img_output_dir}")


# ── CLI ──────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 3:
        print("Uso: python src/slide_images.py <subject> <lesson> [--slides-dir DIR]")
        sys.exit(1)

    subject = sys.argv[1]
    lesson = sys.argv[2]

    sd = None
    if "--slides-dir" in sys.argv:
        idx = sys.argv.index("--slides-dir")
        sd = Path(sys.argv[idx + 1])

    process_sbobina(subject, lesson, slides_dir=sd)


if __name__ == "__main__":
    main()
